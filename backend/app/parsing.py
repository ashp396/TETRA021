import io
import csv
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

def extract_text(filename: str, content: bytes) -> str:
    lower = filename.lower()
    try:
        if lower.endswith(".pdf"):
            return _from_pdf(content)
        if lower.endswith(".docx"):
            return _from_docx(content)
        if lower.endswith(".pptx"):
            return _from_pptx(content)
        if lower.endswith(".xlsx"):
            return _from_xlsx(content)
        if lower.endswith(".csv"):
            return content.decode("utf-8", errors="ignore")
        # fall back to plain text
        return content.decode("utf-8", errors="ignore")
    except Exception as exc:
        raise ValueError(f"Could not read {filename}: {exc}")

def _from_pdf(content: bytes) -> str:
    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    text_parts.append(" | ".join(cell or "" for cell in row))
    return "\n".join(text_parts)

def _from_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)

def _from_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)

def _from_xlsx(content: bytes) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)
